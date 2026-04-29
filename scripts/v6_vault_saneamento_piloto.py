#!/usr/bin/env python3
"""
V6 Vault Saneamento Piloto — Modo Leitura/Sugestão
Analisa arquivos e gera sugestões de renomeação em Markdown.
NÃO executa mv/rename/move.
"""

import os
import re
import sys
import zipfile
from collections import Counter
from typing import Dict, List, Optional, Tuple

BASE_DIR = os.path.expanduser("~/testes/vault/input/ACADEMICO")
REPORT_PATH = os.path.normpath(
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "../reports/V6-ACADEMICO-piloto-renomeacao.md")
)
DEBUG_DIR = os.path.normpath(
    os.path.join(os.path.dirname(os.path.abspath(__file__)), "../debug/V6")
)

SUPPORTED_EXTENSIONS = {'.pdf', '.doc', '.docx', '.odt', '.ppt', '.pptx', '.xls', '.xlsx', '.txt'}

# Extração PDF com OCR (opcional — usa file_extractor.py se disponível)
try:
    from file_extractor import extract_text_with_ocr as _pdf_extract_with_ocr
    from file_extractor import get_ocr_stats as _get_ocr_stats
    HAS_OCR_MODULE = True
except ImportError:
    HAS_OCR_MODULE = False

try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    try:
        import pypdf as PyPDF2  # type: ignore
        HAS_PYPDF2 = True
    except ImportError:
        HAS_PYPDF2 = False

try:
    import docx as python_docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# ---------------------------------------------------------------------------
# Extração de texto
# ---------------------------------------------------------------------------

def clean_text(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r'<[^>]+>', ' ', text)
    text = ''.join(c for c in text if c.isprintable() or c in '\n\r\t')
    lines = []
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if len(line) <= 3:
            lines.append(line)
            continue
        alpha = sum(1 for c in line if c.isalpha())
        if alpha / len(line) > 0.35:
            lines.append(line)
        elif any(kw in line.upper() for kw in [
            'UNIVERSIDADE', 'TITULO', 'TÍTULO', 'AUTOR', 'RESUMO',
            'ABSTRACT', 'DISSERTACAO', 'DISSERTAÇÃO', 'TESE', 'REVISTA',
            'ISSN', 'ARTIGO', 'PUBLICADO', 'DATA', 'ANO',
        ]):
            lines.append(line)
    return '\n'.join(lines)


def extract_pdf_text(filepath: str, max_pages: int = 2, from_end: bool = False) -> str:
    if not HAS_PYPDF2:
        return ""
    try:
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            total = len(reader.pages)
            if total == 0:
                return ""
            pages = range(max(0, total - max_pages), total) if from_end else range(min(max_pages, total))
            parts = []
            for i in pages:
                try:
                    t = reader.pages[i].extract_text()
                    if t:
                        parts.append(clean_text(t))
                except Exception:
                    pass
        return '\n'.join(parts)
    except Exception:
        return ""


def extract_docx_text(filepath: str) -> str:
    if not HAS_DOCX:
        return ""
    try:
        doc = python_docx.Document(filepath)
        return clean_text('\n'.join(p.text for p in doc.paragraphs if p.text.strip()))
    except Exception:
        return ""


def extract_pptx_text(filepath: str) -> str:
    if HAS_PPTX:
        try:
            prs = PptxPresentation(filepath)
            parts = []
            for slide in prs.slides[:10]:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        parts.append(shape.text.strip())
            return clean_text('\n'.join(parts))
        except Exception:
            pass
    try:
        texts = []
        with zipfile.ZipFile(filepath) as z:
            for name in sorted(z.namelist()):
                if re.match(r'ppt/slides/slide\d+\.xml', name):
                    with z.open(name) as f:
                        raw = f.read().decode('utf-8', errors='ignore')
                        texts.extend(re.findall(r'<a:t>([^<]+)</a:t>', raw))
        return clean_text(' '.join(texts[:600]))
    except Exception:
        return ""


def extract_odt_text(filepath: str) -> str:
    try:
        with zipfile.ZipFile(filepath) as z:
            with z.open('content.xml') as f:
                raw = f.read().decode('utf-8', errors='ignore')
        return clean_text(re.sub(r'<[^>]+>', ' ', raw))
    except Exception:
        return ""


def extract_rtf_text(raw_bytes: bytes) -> str:
    try:
        text = raw_bytes.decode('utf-8', errors='ignore')
        while '{\\' in text:
            text = re.sub(r'\{[^{}]*\}', '', text)
        text = re.sub(r'\\[a-z0-9]+', ' ', text)
        return clean_text(re.sub(r'[{}]', '', text))
    except Exception:
        return ""


def extract_text(filepath: str) -> str:
    """Extrai texto de formatos não-PDF. Para PDFs, use _extract_pdf()."""
    ext = os.path.splitext(filepath)[1].lower()

    if ext == '.pdf':
        text = extract_pdf_text(filepath, max_pages=2)
        if len(text) < 100:
            text += extract_pdf_text(filepath, max_pages=2, from_end=True)
        return text

    if ext == '.docx':
        return extract_docx_text(filepath)

    if ext == '.doc':
        try:
            with open(filepath, 'rb') as f:
                raw = f.read(20000)
            text = extract_rtf_text(raw) if raw[:5] == b'{\\rtf' else ""
            if not text:
                for enc in ['utf-8', 'latin-1', 'cp1252']:
                    try:
                        text = clean_text(raw.decode(enc))
                        break
                    except Exception:
                        pass
            return text or clean_text(raw.decode('utf-8', errors='ignore'))
        except Exception:
            return ""

    if ext == '.odt':
        return extract_odt_text(filepath)

    if ext in ('.pptx', '.ppt'):
        return extract_pptx_text(filepath)

    if ext in ('.xlsx', '.xls'):
        try:
            tokens: List[str] = []
            with zipfile.ZipFile(filepath) as z:
                for name in z.namelist():
                    if 'sharedStrings' in name or 'sheet' in name:
                        with z.open(name) as f:
                            raw = f.read().decode('utf-8', errors='ignore')
                            tokens.extend(re.findall(r'<t>([^<]+)</t>', raw))
            return clean_text(' '.join(tokens[:500]))
        except Exception:
            return ""

    if ext == '.txt':
        try:
            with open(filepath, encoding='utf-8', errors='replace') as f:
                return clean_text(f.read(10000))
        except Exception:
            return ""

    return ""


def _extract_pdf(filepath: str) -> 'Tuple[str, bool]':
    """Extrai texto de PDF com fallback OCR se disponível. Retorna (texto, ocr_aplicado)."""
    if HAS_OCR_MODULE:
        return _pdf_extract_with_ocr(filepath)
    text = extract_pdf_text(filepath, max_pages=2)
    if len(text) < 100:
        text += extract_pdf_text(filepath, max_pages=2, from_end=True)
    return text, False


# ---------------------------------------------------------------------------
# Análise do nome do arquivo
# ---------------------------------------------------------------------------

_OPAQUE_RE = [
    re.compile(r'^[\d\-_]+$'),
    re.compile(r'^\d{5,}$'),
    re.compile(r'^[A-Z]{1,3}\d{6,}$'),
    re.compile(r'^[a-z]{2,4}[\-_]\d{4}[\-_]\d+', re.I),
    re.compile(r'^DOC[\-_]\d{8}[\-_]\w+', re.I),
    re.compile(r'^img\d{6,}', re.I),
    re.compile(r'^scan[\-_\d]', re.I),
]

# Tipos de saúde aparecem ANTES dos tipos genéricos para evitar sobreposição
_NAME_TYPE_MAP = [
    # Saúde — tipos específicos (verificados antes dos genéricos)
    # Sem \b final para capturar nomes compostos (LaudoBiopsia001, resultadoexames, etc.)
    (r'\bresultado',                            'Resultado de Exame'),
    (r'\blaudo',                               'Laudo'),
    (r'\bhemograma',                           'Resultado de Exame'),
    (r'\bultrassom|\bultrassonografia',        'Resultado de Exame'),
    (r'\bresson[aâ]ncia',                      'Resultado de Exame'),
    (r'\bcintilografia',                       'Resultado de Exame'),
    (r'\btomografia',                          'Resultado de Exame'),
    (r'\bbiopsia|\bbiópsia',                   'Laudo'),
    (r'\bquestion[aá]rio',                     'Questionário Clínico'),
    (r'\breceitu[aá]rio',                      'Receita'),
    (r'\bpreparo\s+(?:de|para)',               'Guia de Exame'),
    (r'\borientac[aã]o\s+(?:de|para)',         'Guia de Exame'),
    # Acadêmico e documentos gerais
    (r'\baula\b',           'Aula'),
    (r'\bapresentac[aã]o',  'Apresentação'),
    (r'\bartigo\b',         'Artigo Científico'),
    (r'\bdissertac[aã]o',   'Dissertação'),
    (r'\btese\b',           'Tese'),
    (r'\bmonografia',       'Monografia'),
    (r'\btcc\b',            'TCC'),
    (r'\btrabalho',         'Trabalho Acadêmico'),
    # Financeiro — tipos específicos (antes dos genéricos)
    (r'\birpf\b',                             'Declaração de IRPF'),
    (r'\bextrato\b',                          'Extrato Bancário'),
    (r'\bholerite\b|\bcontracheque\b',        'Recibo de Salário/Holerite'),
    (r'nf-?e\b|\bdanfe\b',                    'Nota Fiscal'),
    (r'\bfatura\b',                           'Fatura'),
    (r'\bboleto',           'Boleto'),
    (r'\bcomprovante',      'Comprovante'),
    (r'\bdiploma',          'Diploma'),
    (r'\bcertificado',      'Certificado'),
    (r'\bcurr[ií]culo',     'Currículo'),
    (r'\brelat[oó]rio',     'Relatório'),
    (r'\bresumo\b',         'Resumo'),
    (r'\bexerc[ií]cio',     'Exercício'),
    (r'\bavalia[cç][aã]o',  'Avaliação'),
    (r'\bedital\b',         'Edital'),
    (r'\bprojeto\b',        'Projeto'),
    (r'\bcarta\b',          'Carta'),
    (r'\bdeclara[cç][aã]o', 'Declaração'),
    (r'\broteiro\b',        'Roteiro'),
    (r'\bguia\b',           'Guia'),
    (r'\bplanilha',         'Planilha'),
    (r'\bprova\b',          'Prova'),
    (r'\bgabarito',         'Gabarito'),
    (r'\bapostila',         'Apostila'),
    (r'\bmanual\b',         'Manual'),
    (r'\bformul[aá]rio',    'Formulário'),
    (r'\bficha\b',          'Ficha'),
    (r'\bcontrato',         'Contrato'),
    (r'\bminuta\b',         'Minuta'),
    (r'\brequerimento',     'Requerimento'),
    (r'hist[oó]rico',       'Histórico'),
    (r'\bprint',            'Screenshot'),
]


def analyze_filename(filename: str) -> Dict:
    stem, ext = os.path.splitext(filename)
    stem = stem.strip()
    stem_lower = stem.lower()

    is_opaque = any(p.match(stem) for p in _OPAQUE_RE)

    if not is_opaque:
        alpha_words = re.findall(r'[a-záéíóúàâêôãõçüA-ZÁÉÍÓÚÀÂÊÔÃÕÇÜ]{3,}', stem)
        if len(alpha_words) < 2:
            alpha_ratio = sum(1 for c in stem if c.isalpha()) / max(len(stem), 1)
            if alpha_ratio < 0.40:
                is_opaque = True

    tipo_from_name: Optional[str] = None
    for pattern, tipo in _NAME_TYPE_MAP:
        if re.search(pattern, stem_lower):
            tipo_from_name = tipo
            break

    if is_opaque and tipo_from_name:
        alpha_words = re.findall(r'[a-záéíóúàâêôãõçüA-ZÁÉÍÓÚÀÂÊÔÃÕÇÜ]{3,}', stem)
        if len(alpha_words) >= 3:
            is_opaque = False

    return {
        'stem': stem,
        'ext': ext,
        'is_opaque': is_opaque,
        'tipo_from_name': tipo_from_name,
    }


# ---------------------------------------------------------------------------
# Detector especializado: artigo de periódico
# ---------------------------------------------------------------------------

_INST_WORD_KW = frozenset([
    'UNIVERSIDADE', 'FACULDADE', 'INSTITUTO', 'DEPARTAMENTO',
    'PROGRAMA', 'CURSO', 'CENTRO', 'ESCOLA', 'DISCIPLINA',
    'COORDENAÇÃO', 'COORDENACAO', 'GRADUAÇÃO', 'GRADUACAO',
    'MESTRADO', 'DOUTORADO', 'FUNDAÇÃO', 'FUNDACAO',
    'RESUMO', 'ABSTRACT', 'PALAVRAS', 'KEYWORDS', 'TÍTULO', 'TITULO',
])

_NOMES_RELEVANTES_LIST = [
    ('Conrado',          re.compile(r'\bconrado\b',          re.I)),
    ('Nogueira',         re.compile(r'\bnogueira\b',         re.I)),
    ('Emilia',           re.compile(r'\bemíli[ao]\b|\bemili[ao]\b', re.I)),
    ('Sonia',            re.compile(r'\bsônia\b|\bsonia\b',  re.I)),
    ('Rafael Cavalcante',re.compile(r'\brafael\s+cavalcante\b', re.I)),
    ('Norte Certo',      re.compile(r'\bnorte\s+certo\b',    re.I)),
    ('North Paradise',   re.compile(r'\bnorth\s+paradise\b', re.I)),
    ('Linguiçaria',      re.compile(r'\blinguiçaria\b|\blinguicaria\b', re.I)),
    ('Eu Gourmet',       re.compile(r'\beu\s+gourmet\b',     re.I)),
    ('Cheffon',          re.compile(r'\bcheffon\b',          re.I)),
    ('Rock City',        re.compile(r'\brock\s+city\b',      re.I)),
    ('Casa Lisboa',      re.compile(r'\bcasa\s+lisboa\b',    re.I)),
    ('Reserva 19',       re.compile(r'\breserva\s+19\b',     re.I)),
]

_PERIODICO_RE = re.compile(
    r'([A-ZÁÉÍÓÚ][^,\n]{4,70}),\s*\d+\s*\(\d+\)\s*,\s*p\.\s*\d+',
    re.IGNORECASE,
)
_PERIODICO_ANO_RE = re.compile(
    r'([A-ZÁÉÍÓÚ][^,\n]{4,70}),\s*\d+\s*\(\d+\)\s*[^,\n]*,\s*((?:19|20)\d{2})',
    re.IGNORECASE,
)


def detecta_artigo_academico(text: str) -> Optional[Dict]:
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if not lines:
        return None

    periodico = ""
    ano_periodico = ""
    for line in lines[:20]:
        m_ano = _PERIODICO_ANO_RE.search(line)
        if m_ano:
            periodico = m_ano.group(1).strip()
            ano_periodico = m_ano.group(2)
            break
        m = _PERIODICO_RE.search(line)
        if m:
            periodico = m.group(1).strip()
            break

    has_issn = any('ISSN' in l.upper() for l in lines[:30])

    if not periodico and not has_issn:
        return None

    titulo = ""
    TITLE_EXCL = {
        'UNIVERSIDADE', 'FACULDADE', 'INSTITUTO', 'DEPARTAMENTO',
        'PROGRAMA', 'CURSO', 'CENTRO', 'FUNDAÇÃO', 'FUNDACAO',
        'REVISTA', 'ISSN', 'WWW', 'HTTP',
    }
    for i, line in enumerate(lines[:80]):
        upper = line.upper()
        if upper.startswith(('RESUMO', 'ABSTRACT', 'PALAVRAS', 'KEYWORDS')):
            candidates = []
            for prev in lines[max(0, i - 15):i]:
                stripped = re.sub(r'^\d+\s+', '', prev).strip()
                if (15 <= len(stripped) <= 250
                        and not any(kw in stripped.upper() for kw in TITLE_EXCL)
                        and not stripped[0].isdigit()
                        and not _PERIODICO_RE.search(prev)):
                    candidates.append(stripped)
            if candidates:
                titulo = max(candidates, key=len)[:200]
            break

    if not titulo:
        for line in lines[:30]:
            stripped = re.sub(r'^\d+\s+', '', line).strip()
            if (20 <= len(stripped) <= 200
                    and not any(kw in stripped.upper() for kw in TITLE_EXCL)
                    and not stripped[0].isdigit()
                    and not _PERIODICO_RE.search(line)):
                titulo = stripped[:200]
                break

    autores: List[str] = []
    INST_KW_SET = {
        'UNIVERSIDADE', 'FACULDADE', 'INSTITUTO', 'DEPARTAMENTO', 'PROGRAMA',
        'CURSO', 'CENTRO', 'FUNDAÇÃO', 'BRASIL', 'PARÁ', 'AMAPÁ', 'RESUMO',
        'ABSTRACT', 'PALAVRAS', 'KEYWORDS',
    }
    for line in lines[:60]:
        if re.search(r'[,;():!?]', line):
            continue
        line_clean = re.sub(r'\s+\d+\s*$', '', line).strip()
        words = line_clean.split()
        if 2 <= len(words) <= 5:
            if any(len(w) > 12 for w in words):
                continue
            if all(re.match(r'^[a-záéíóúàâêôãõçüA-ZÁÉÍÓÚÀÂÊÔÃÕÇÜ]+$', w) for w in words):
                if words[0][0].isupper():
                    if not any(kw in line_clean.upper() for kw in INST_KW_SET):
                        autores.append(line_clean)
                        if len(autores) >= 3:
                            break

    sobrenome = ""
    if autores:
        parts = autores[0].split()
        sobrenome = parts[-1] if parts else ""

    ano = ano_periodico
    if not ano:
        m = re.search(r',\s*((?:19|20)\d{2})\s*[-–]', text)
        if m:
            ano = m.group(1)
    if not ano:
        years = [int(y) for y in re.findall(r'\b(19\d{2}|20[0-2]\d)\b', text)]
        plausible = [y for y in years if 1980 <= y <= 2030]
        if plausible:
            ano = str(Counter(plausible).most_common(1)[0][0])

    return {
        'tipo_documento': 'Artigo Científico',
        'titulo': titulo,
        'autores': autores,
        'autor_principal_sobrenome': sobrenome,
        'ano': ano,
        'periodico': periodico,
    }


# ---------------------------------------------------------------------------
# Detector especializado: capa de trabalho acadêmico
# ---------------------------------------------------------------------------

def detecta_capa_academica(text: str) -> Optional[Dict]:
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if len(lines) < 4:
        return None

    if not any(any(kw in line.upper() for kw in _INST_WORD_KW) for line in lines[:40]):
        return None

    nome_idx = -1
    autor_principal = ""
    for i, line in enumerate(lines[:80]):
        if re.search(r'[,;():!?/\\0-9]', line):
            continue
        words = line.split()
        if not (2 <= len(words) <= 5):
            continue
        if not all(re.match(r'^[a-záéíóúàâêôãõçüA-ZÁÉÍÓÚÀÂÊÔÃÕÇÜ]+$', w) for w in words):
            continue
        if not words[0][0].isupper():
            continue
        if any(kw in line.upper() for kw in _INST_WORD_KW):
            continue
        if not any(len(w) >= 3 for w in words):
            continue
        nome_idx = i
        autor_principal = line
        break

    if nome_idx == -1 or not autor_principal:
        return None

    titulo = ""
    for line in lines[nome_idx + 1: nome_idx + 7]:
        if len(line) >= 10 and not any(kw in line.upper() for kw in _INST_WORD_KW):
            titulo = line[:200]
            break

    if not titulo:
        return None

    if autor_principal == autor_principal.upper():
        autor_principal = autor_principal.title()

    years = [int(y) for y in re.findall(r'\b(19\d{2}|20[0-2]\d)\b', text)]
    plausible = [y for y in years if 1980 <= y <= 2030]
    ano = str(Counter(plausible).most_common(1)[0][0]) if plausible else ""

    return {
        'tipo_documento': 'Trabalho Acadêmico',
        'titulo': titulo,
        'autor_principal': autor_principal,
        'ano': ano,
    }


# ---------------------------------------------------------------------------
# Detector de nomes/termos relevantes
# ---------------------------------------------------------------------------

def detecta_nomes_relevantes(text: str) -> List[str]:
    return [nome for nome, pat in _NOMES_RELEVANTES_LIST if pat.search(text)]


# ---------------------------------------------------------------------------
# Detecção de documentos de saúde — short-circuit para pasta /SAUDE/
# ---------------------------------------------------------------------------

# Palavras-chave de resultado de exame laboratorial
_EXAME_LAB_KW = frozenset([
    'RESULTADO DE EXAME', 'EXAMES DE ANÁLISES', 'EXAME DE SANGUE',
    'HEMOGRAMA', 'HEMATÓCRITO', 'HEMATOCRITO', 'LEUCÓCITOS', 'LEUCOCITOS',
    'PLAQUETAS', 'COAGULOGRAMA', 'GLICOSE', 'COLESTEROL', 'TRIGLICERÍDEOS',
    'TRIGLICERIDEOS', 'ANÁLISES CLÍNICAS', 'ANALISES CLINICAS',
    'LABORATÓRIO', 'LABORATORIO', 'DIAGNOSE', 'RESULTADO:',
])

# Palavras-chave de exames de imagem / procedimento
_EXAME_IMAGEM_KW = frozenset([
    'RESSONÂNCIA MAGNÉTICA', 'RESSONANCIA MAGNETICA',
    'TOMOGRAFIA COMPUTADORIZADA', 'ULTRASSONOGRAFIA', 'ULTRASSOM',
    'CINTILOGRAFIA', 'RADIOGRAFIA', 'MAMOGRAFIA', 'ECOCARDIOGRAMA',
    'ELETROCARDIOGRAMA', 'DENSITOMETRIA', 'ENDOSCOPIA', 'COLONOSCOPIA',
])

# Palavras-chave de laudo / patologia
_LAUDO_KW = frozenset([
    'LAUDO', 'HISTOPATOLÓGICO', 'HISTOPATOLOGICO',
    'ANATOMIA PATOLÓGICA', 'ANATOMIA PATOLOGICA',
    'BIÓPSIA', 'BIOPSIA', 'CITOLOGIA', 'IMUNOHISTOQUÍMICA',
    'PARECER MÉDICO', 'PARECER PSICOLÓGICO',
])

# Palavras-chave de guia / orientação de exame
_GUIA_EXAME_KW = frozenset([
    'PREPARO PARA', 'PREPARAÇÃO PARA', 'PREPARACAO PARA',
    'ORIENTAÇÕES PARA', 'ORIENTACOES PARA',
    'INSTRUÇÕES PARA', 'INSTRUCOES PARA',
    'ANTES DO EXAME', 'JEJUM DE',
])

# Palavras-chave de questionário clínico
_QUESTIONARIO_CLINICO_KW = frozenset([
    'QUESTIONÁRIO', 'QUESTIONARIO',
    'INVENTÁRIO', 'INVENTARIO',
    'ESCALA DE', 'YOUNG', 'BECK', 'HAMILTON',
    'PHQ-', 'GAD-',
    'VULNERABILIDADE AO', 'ESQUEMAS INICIAIS',
])

# Palavras-chave de receita / prescrição
_RECEITA_KW = frozenset([
    'RECEITUÁRIO', 'RECEITUARIO', 'RECEITA MÉDICA', 'RECEITA MEDICA',
    'POSOLOGIA', 'PRESCRIÇÃO', 'PRESCRICAO',
])

# Sinais fracos de contexto de saúde
_SAUDE_FRACO_KW = frozenset([
    'UNIMED', 'PACIENTE', 'CRM', 'CRF', 'CRP',
    'PRONTUÁRIO', 'PRONTUARIO', 'MÉDICO', 'MEDICO',
    'CLÍNICA', 'CLINICA', 'HOSPITAL',
])

# Conjunto de todos os tipos considerados documentos de saúde
_SAUDE_TIPOS = frozenset([
    'Resultado de Exame', 'Laudo', 'Guia de Exame', 'Questionário Clínico',
    'Receita', 'Documento de Saúde', 'Orientação de Exame', 'Parecer Psicológico',
])


# ---------------------------------------------------------------------------
# Detecção de documentos financeiros — short-circuit para pasta /FINANCEIRO/
# ---------------------------------------------------------------------------

_NFE_KW = frozenset([
    'NOTA FISCAL', 'NF-E', 'NFE', 'DANFE', 'CHAVE DE ACESSO',
    'DOCUMENTO AUXILIAR DA NOTA FISCAL ELETRÔNICA',
    'DOCUMENTO AUXILIAR DA NOTA FISCAL ELETRONICA',
])

_IRPF_KW = frozenset([
    'IRPF', 'IMPOSTO DE RENDA', 'DECLARAÇÃO DE AJUSTE ANUAL',
    'DECLARACAO DE AJUSTE ANUAL', 'RECEITA FEDERAL',
    'DECLARAÇÃO DE IMPOSTO DE RENDA', 'DECLARACAO DE IMPOSTO DE RENDA',
])

_EXTRATO_KW = frozenset([
    'EXTRATO BANCÁRIO', 'EXTRATO BANCARIO', 'EXTRATO DE CONTA',
    'EXTRATO MENSAL', 'MOVIMENTAÇÃO FINANCEIRA', 'MOVIMENTACAO FINANCEIRA',
    'SALDO ANTERIOR', 'SALDO ATUAL', 'CONTA CORRENTE',
])

_HOLERITE_KW = frozenset([
    'HOLERITE', 'CONTRACHEQUE', 'RECIBO DE SALÁRIO', 'RECIBO DE SALARIO',
    'FOLHA DE PAGAMENTO', 'SALÁRIO BRUTO', 'SALARIO BRUTO',
    'SALÁRIO LÍQUIDO', 'SALARIO LIQUIDO', 'INSS', 'FGTS',
    'VALE TRANSPORTE', 'VALE REFEIÇÃO',
])

_COMPROVANTE_TRANSF_KW = frozenset([
    'COMPROVANTE DE TRANSFERÊNCIA', 'COMPROVANTE DE TRANSFERENCIA',
    'TRANSFERÊNCIA PIX', 'TRANSFERENCIA PIX', 'PIX ENVIADO',
    'TED ENVIADO', 'DOC BANCÁRIO', 'DOC BANCARIO',
    'VALOR TRANSFERIDO', 'CONTA DE DESTINO', 'CONTA DE ORIGEM',
])

_FATURA_KW = frozenset([
    'FATURA DO CARTÃO', 'FATURA DO CARTAO',
    'FATURA DE CARTÃO DE CRÉDITO', 'FATURA DE CARTAO DE CREDITO',
    'TOTAL DA FATURA', 'PAGAMENTO MÍNIMO', 'PAGAMENTO MINIMO',
    'LIMITE DE CRÉDITO', 'LIMITE DE CREDITO',
])

_DARF_KW = frozenset([
    'DARF', 'DOCUMENTO DE ARRECADAÇÃO', 'DOCUMENTO DE ARRECADACAO',
    'RECEITA FEDERAL', 'CÓDIGO DE RECEITA', 'CODIGO DE RECEITA',
])

_BOLETO_KW = frozenset([
    'BOLETO BANCÁRIO', 'BOLETO BANCARIO', 'NOSSO NÚMERO', 'NOSSO NUMERO',
    'LINHA DIGITÁVEL', 'LINHA DIGITAVEL', 'CÓDIGO DE BARRAS', 'CODIGO DE BARRAS',
    'BENEFICIÁRIO', 'BENEFICIARIO',
])

_FINANCEIRO_FRACO_KW = frozenset([
    'BANCO', 'AGÊNCIA', 'AGENCIA', 'CPF', 'CNPJ',
    'PAGAMENTO', 'RECEBIMENTO', 'CRÉDITO', 'DÉBITO',
])

_FINANCEIRO_TIPOS = frozenset([
    'Nota Fiscal', 'Comprovante de Transferência', 'Recibo de Salário/Holerite',
    'Extrato Bancário', 'Declaração de IRPF', 'DARF', 'Fatura', 'Boleto',
    'Documento Financeiro',
])


def _has_any(text_upper: str, kw_set: frozenset) -> bool:
    return any(kw in text_upper for kw in kw_set)


def _clean_metadata_string(s: str, max_len: int = 60) -> str:
    """Limpa e trunca strings de metadados para uso no nome do arquivo."""
    if not s:
        return ""
    # Remove excesso de espaços e quebras de linha
    s = re.sub(r'\s+', ' ', s).strip()
    
    # Filtra padrões claramente indesejados (ruído de OCR/extração)
    bad_tokens = [
        r'FORMA FARMACÊUTICA', r'CNPJ FOL', r'CNPJ/CPF/IdEstrangeiro',
        r'VENCIMENTO', r'PAGAR', r'VALOR', r'TOTAL', r'PAGAMENTO',
        r'LINHA DIGITÁVEL', r'CÓDIGO DE BARRAS', r'DOCUMENTO',
        r'CONTRAIU DESPESAS', r'OPERACIONAIS', r'DESPESAS MÉDICAS',
    ]
    for pat in bad_tokens:
        if re.search(pat, s, re.IGNORECASE):
            # Se a string contém muito lixo, tentamos reduzir ou descartar
            parts = re.split(pat, s, flags=re.IGNORECASE)
            s = parts[0].strip()
            if len(s) < 3:
                return ""
            break
            
    # Remove caracteres especiais não amigáveis a nomes de arquivo
    s = re.sub(r'[/\\:*?"<>|]', '-', s)
    
    # Remove repetições estranhas (ex: "Conrado Conrado")
    words = s.split()
    if len(words) >= 2 and words[0] == words[1]:
        s = " ".join(words[1:])

    if len(s) > max_len:
        s = s[:max_len].strip() + "…"
    return s.strip()


def _saude_score(text_upper: str) -> int:
    """Pontua evidências de saúde no texto. Score >= 3 → contexto de saúde forte."""
    score = 0
    if _has_any(text_upper, _EXAME_LAB_KW):
        score += 3
    if _has_any(text_upper, _EXAME_IMAGEM_KW):
        score += 2
    if _has_any(text_upper, _LAUDO_KW):
        score += 3
    if _has_any(text_upper, _GUIA_EXAME_KW):
        score += 2
    if _has_any(text_upper, _QUESTIONARIO_CLINICO_KW):
        score += 2
    if _has_any(text_upper, _RECEITA_KW):
        score += 2
    if _has_any(text_upper, _SAUDE_FRACO_KW):
        score += 1
    return score


def _in_saude_path(filepath: str) -> bool:
    return '/SAUDE/' in filepath.replace('\\', '/').upper()


def _is_saude_evidence(text: str) -> bool:
    """True se o texto tem evidência forte de documento de saúde (fora da pasta /SAUDE/)."""
    if not text:
        return False
    return _saude_score(text.upper()) >= 3


def _in_financeiro_path(filepath: str) -> bool:
    return '/FINANCEIRO/' in filepath.replace('\\', '/').upper()


def _financeiro_score(text_upper: str) -> int:
    """Pontua evidências financeiras no texto. Score >= 3 → contexto financeiro forte."""
    score = 0
    if _has_any(text_upper, _NFE_KW):
        score += 4
    if _has_any(text_upper, _IRPF_KW):
        score += 3
    if _has_any(text_upper, _EXTRATO_KW):
        score += 3
    if _has_any(text_upper, _HOLERITE_KW):
        score += 4
    if _has_any(text_upper, _COMPROVANTE_TRANSF_KW):
        score += 3
    if _has_any(text_upper, _FATURA_KW):
        score += 2
    if _has_any(text_upper, _BOLETO_KW):
        score += 3
    if _has_any(text_upper, _FINANCEIRO_FRACO_KW):
        score += 1
    return score


def _is_financeiro_evidence(text: str) -> bool:
    """True se o texto tem evidência forte de documento financeiro (fora de /FINANCEIRO/)."""
    if not text:
        return False
    return _financeiro_score(text.upper()) >= 3


def detecta_documento_financeiro(text: str, filepath: str) -> Optional[Dict]:
    """
    Detecta tipo de documento financeiro e extrai metadados básicos.
    Retorna None se não for possível determinar um tipo específico financeiro.
    """
    u = text.upper() if text else ""
    fname_lower = os.path.basename(filepath).lower()

    tipo: Optional[str] = None
    
    # Operadoras/Serviços comuns que indicam Fatura
    OPERADORAS = ['NET', 'CLARO', 'VIVO', 'OI', 'TIM', 'COELBA', 'CPFL', 'ENERGIA', 'ÁGUA', 'AGUA', 'SANEAMENTO', 'NUBANK', 'ITAU', 'BRADESCO', 'SANTANDER', 'CARTÃO', 'CARTAO']
    has_op_in_name = any(op in fname_lower.upper() for op in OPERADORAS)

    # Ordem de prioridade: DARF > IRPF > NF-e > Holerite > Extrato > Comprovante > Fatura > Boleto > Fraco
    if _has_any(u, _DARF_KW) or re.search(r'darf|documento.?arrecada[cç][aã]o', fname_lower):
        tipo = 'DARF'
    elif _has_any(u, _IRPF_KW) or re.search(r'irpf|imposto.?renda', fname_lower):
        tipo = 'Declaração de IRPF'
    elif _has_any(u, _NFE_KW) or re.search(r'nf-?e|danfe|nota.?fiscal', fname_lower):
        tipo = 'Nota Fiscal'
    elif _has_any(u, _HOLERITE_KW) or re.search(r'holerite|contracheque', fname_lower):
        tipo = 'Recibo de Salário/Holerite'
    elif _has_any(u, _EXTRATO_KW) or re.search(r'extrato', fname_lower):
        tipo = 'Extrato Bancário'
    elif _has_any(u, _COMPROVANTE_TRANSF_KW) or re.search(
        r'comprovante|transfer[eê]ncia|pix', fname_lower
    ):
        tipo = 'Comprovante de Transferência'
    elif _has_any(u, _FATURA_KW) or re.search(r'fatura', fname_lower) or (has_op_in_name and re.search(r'boleto', fname_lower)):
        tipo = 'Fatura'
    elif _has_any(u, _BOLETO_KW) or re.search(r'boleto', fname_lower):
        tipo = 'Boleto'
    elif _has_any(u, _FINANCEIRO_FRACO_KW):
        tipo = 'Documento Financeiro'
    elif _in_financeiro_path(filepath):
        tipo = 'Documento Financeiro'
    else:
        return None

    lines = [l.strip() for l in text.split('\n') if l.strip()] if text else []

    # Emitente / Beneficiário / Titular
    entity = ""
    for line in lines[:40]:
        m = re.search(
            r'(?:emitente|benefici[aá]rio|raz[aã]o\s+social|empresa|contribuinte|titular|cliente|assinante|nome|pagador)[:\s]+([A-ZÁÉÍÓÚÀÂÊÔÃÕÇ][^\n]{4,60})',
            line, re.IGNORECASE,
        )
        if m:
            candidate = m.group(1).strip().rstrip('.,:;')
            if 4 <= len(candidate) <= 80 and not re.search(r'[();]', candidate):
                entity = candidate[:80]
                break

    # Prioridade absoluta para nomes relevantes conhecidos (evita erros de OCR como CONARDO)
    if text:
        for nome, pat in _NOMES_RELEVANTES_LIST:
            if pat.search(text):
                entity = nome
                break

    # Título ou Operadora (especialmente para Fatura)
    titulo = ""
    if tipo == 'Fatura':
        for op in OPERADORAS:
            if op in u or op in fname_lower.upper():
                titulo = op.capitalize()
                break
    
    if not titulo:
        FIN_INST = {'BANCO', 'AGÊNCIA', 'AGENCIA', 'CNPJ', 'CPF', 'CNPJ/CPF', 'VENCIMENTO', 'PAGAMENTO'}
        for line in lines[:10]:
            if len(line) >= 8 and not any(kw in line.upper() for kw in FIN_INST):
                titulo = line[:150]
                break

    # Data: tenta encontrar Mês/Ano (AAAAMM) para faturas ou apenas Ano
    ano = ""
    # Para Faturas e comprovantes mensais, prioriza Mês/Ano (dd/mm/aaaa -> dd, mm, aaaa)
    if tipo in ('Fatura', 'Extrato Bancário', 'Recibo de Salário/Holerite'):
        # Procura por dd/mm/aaaa ou mm/aaaa
        m_data = re.search(r'\b(?:(?:\d{2}/)?(0[1-9]|1[0-2])/(\d{4}))\b', text) if text else None
        if m_data:
            mm, aaaa = m_data.groups()
            if 1980 <= int(aaaa) <= 2030:
                ano = f"{aaaa}-{mm}"

    if not ano:
        years = [int(y) for y in re.findall(r'\b(19\d{2}|20[0-2]\d)\b', text)] if text else []
        plausible = [y for y in years if 1980 <= y <= 2030]
        if plausible:
            counts = Counter(plausible)
            if tipo == 'Declaração de IRPF':
                # No IRPF, tenta detectar Exercício e Base
                top_years = sorted(list(set(plausible)), reverse=True)
                if len(top_years) >= 2:
                    y1, y2 = top_years[0], top_years[1]
                    # IRPF costuma ter o ano do exercício e o ano-calendário (ex: 2018 e 2017)
                    if y1 == y2 + 1:
                        ano = f"{y1}-Base{y2}"
                    else:
                        ano = str(y1)
                else:
                    ano = str(top_years[0])
            else:
                ano = str(counts.most_common(1)[0][0])

    return {
        'tipo_documento': tipo,
        'titulo': titulo,
        'autor_principal': entity,
        'ano': ano,
    }


def detecta_documento_saude(text: str, filepath: str) -> Optional[Dict]:
    """
    Detecta tipo de documento de saúde e extrai metadados básicos.
    Retorna None se não for possível determinar um tipo específico de saúde.
    """
    u = text.upper() if text else ""
    fname_lower = os.path.basename(filepath).lower()

    tipo: Optional[str] = None

    # Ordem de prioridade: Questionário > Laudo > Resultado > Guia > Receita > Imagem > Fraco
    if _has_any(u, _QUESTIONARIO_CLINICO_KW) or re.search(r'question[aá]rio', fname_lower):
        tipo = 'Questionário Clínico'
    elif _has_any(u, _LAUDO_KW) or re.search(r'laudo|biopsia|biópsia', fname_lower):
        tipo = 'Laudo'
    elif _has_any(u, _EXAME_LAB_KW) or re.search(r'resultado|hemograma', fname_lower):
        tipo = 'Resultado de Exame'
    elif _has_any(u, _GUIA_EXAME_KW) or re.search(r'preparo|orientac', fname_lower):
        tipo = 'Guia de Exame'
    elif _has_any(u, _RECEITA_KW) or re.search(r'receitu', fname_lower):
        tipo = 'Receita'
    elif _has_any(u, _EXAME_IMAGEM_KW) or re.search(
        r'ressonancia|tomografia|ultrassom|cintilografia|radiografia|mamografia', fname_lower
    ):
        tipo = 'Resultado de Exame'
    elif _has_any(u, _SAUDE_FRACO_KW):
        tipo = 'Documento de Saúde'
    elif _in_saude_path(filepath):
        tipo = 'Documento de Saúde'
    else:
        return None

    lines = [l.strip() for l in text.split('\n') if l.strip()] if text else []

    # Paciente: busca explícita por campo "paciente"
    patient = ""
    for line in lines[:40]:
        m = re.search(
            r'(?:paciente|patient|nome)[:\s]+([A-ZÁÉÍÓÚÀÂÊÔÃÕÇ][^\n]{4,60})',
            line, re.IGNORECASE,
        )
        if m:
            candidate = m.group(1).strip().rstrip('.,:;')
            if 4 <= len(candidate) <= 80 and not re.search(r'[();]', candidate):
                patient = candidate[:80]
                break

    # Fallback: nomes relevantes conhecidos no texto
    if not patient and text:
        for nome, pat in _NOMES_RELEVANTES_LIST:
            if pat.search(text):
                # tenta extrair o nome completo capitalizado ao redor do match
                m2 = pat.search(text)
                if m2:
                    ctx = text[max(0, m2.start() - 30): m2.end() + 60]
                    mc = re.search(
                        r'([A-ZÁÉÍÓÚÀÂÊÔÃÕÇ][A-Za-záéíóúàâêôãõç ]{5,50})',
                        ctx,
                    )
                    if mc:
                        patient = mc.group(1).strip()
                break

    # Título: primeira linha significativa sem termos institucionais de saúde
    SAUDE_INST = {'UNIMED', 'LABORATÓRIO', 'LABORATORIO', 'CLÍNICA', 'CLINICA',
                  'HOSPITAL', 'MÉDICO', 'MEDICO', 'DIAGNOSE'}
    titulo = ""
    for line in lines[:15]:
        if len(line) >= 10 and not any(kw in line.upper() for kw in SAUDE_INST):
            titulo = line[:150]
            break

    # Ano: mais frequente no texto
    years = [int(y) for y in re.findall(r'\b(19\d{2}|20[0-2]\d)\b', text)] if text else []
    plausible = [y for y in years if 1980 <= y <= 2030]
    ano = str(Counter(plausible).most_common(1)[0][0]) if plausible else ""

    return {
        'tipo_documento': tipo,
        'titulo': titulo,
        'autor_principal': patient,
        'ano': ano,
    }


# ---------------------------------------------------------------------------
# Extração genérica de metadados
# ---------------------------------------------------------------------------

def detect_type_from_content(text: str) -> Optional[str]:
    u = text.upper()
    if any(k in u for k in ['DISSERTAÇÃO APRESENTADA', 'DISSERTACAO APRESENTADA', 'DISSERTAÇÃO DE MESTRADO']):
        return 'Dissertação'
    if any(k in u for k in ['TESE DE DOUTORADO', 'TESE APRESENTADA']):
        return 'Tese'
    if 'ISSN' in u:
        return 'Artigo Científico'
    if any(k in u for k in ['RESUMO', 'ABSTRACT']) and any(k in u for k in ['PALAVRAS-CHAVE', 'KEYWORDS']):
        return 'Artigo Científico'
    if 'DIPLOMA' in u:
        return 'Diploma'
    if 'CERTIFICADO' in u and 'CERTIFICA' in u:
        return 'Certificado'
    if any(k in u for k in ['BOLETO BANCÁRIO', 'BOLETO BANCARIO', 'NOSSO NÚMERO', 'NOSSO NUMERO']):
        return 'Boleto'
    if any(k in u for k in ['COMPROVANTE DE MATRÍCULA', 'COMPROVANTE DE MATRICULA']):
        return 'Comprovante de Matrícula'
    if 'REQUERIMENTO' in u:
        return 'Requerimento'
    if 'EDITAL' in u:
        return 'Edital'
    return None


def extract_title_from_content(text: str, doc_type: str) -> str:
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if not lines:
        return ""
    INST_KW = {'UNIVERSIDADE', 'FACULDADE', 'INSTITUTO', 'DEPARTAMENTO',
                'PROGRAMA', 'CURSO', 'CENTRO', 'FUNDAÇÃO', 'FUNDACAO'}
    if doc_type in ('Artigo Científico', 'Dissertação', 'Tese', 'Monografia', 'TCC'):
        for i, line in enumerate(lines[:60]):
            if line.upper().startswith(('RESUMO', 'ABSTRACT', 'PALAVRAS')):
                candidates = [
                    l for l in lines[max(0, i - 12):i]
                    if 15 <= len(l) <= 200
                    and not any(kw in l.upper() for kw in INST_KW)
                    and not l[0].isdigit()
                ]
                if candidates:
                    return max(candidates, key=len)[:180]
                break
    for line in lines[:25]:
        if 20 <= len(line) <= 200 and not any(kw in line.upper() for kw in INST_KW):
            return line[:180]
    return ""


def extract_author_from_content(text: str) -> str:
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    for line in lines[:80]:
        m = re.search(
            r'(?:Autor[ae]?s?|Alun[oa]|Discente)[\s:]+([A-ZÁÉÍÓÚÀÂÊÔÃÕÇ][^\n]{4,60})',
            line,
        )
        if m:
            candidate = m.group(1).strip()
            if not re.search(r'[,;():!?]', candidate):
                return candidate[:100]

    INST_KW = {'UNIVERSIDADE', 'FACULDADE', 'INSTITUTO', 'DEPARTAMENTO',
                'PROGRAMA', 'CURSO', 'CENTRO', 'FUNDAÇÃO', 'BRASIL',
                'PARÁ', 'AMAPÁ', 'TÍTULO', 'TRABALHO', 'RESUMO'}
    for line in lines[3:40]:
        if re.search(r'[,;():!?]', line):
            continue
        words = line.split()
        if 2 <= len(words) <= 5:
            if all(re.match(r'^[a-záéíóúàâêôãõçüA-ZÁÉÍÓÚÀÂÊÔÃÕÇÜ]+$', w) for w in words):
                if words[0][0].isupper():
                    if not any(kw in line.upper() for kw in INST_KW):
                        return line[:100]
    return ""


def extract_year_from_content(text: str) -> str:
    contextual = [
        r'(?:publicado|publicação|emitido|emissão|data)[^\n]{0,40}((?:19|20)\d{2})',
        r'[©\(c\)]\s*((?:19|20)\d{2})',
        r'ISSN[^\n]{5,80}((?:19|20)\d{2})',
        r',\s*((?:19|20)\d{2})\s*[.,\)]',
        r'\b((?:19|20)\d{2})\b\s*$',
    ]
    for pat in contextual:
        m = re.search(pat, text, re.IGNORECASE | re.MULTILINE)
        if m:
            y = int(m.group(1))
            if 1980 <= y <= 2030:
                return str(y)
    years = [int(y) for y in re.findall(r'\b(19\d{2}|20[0-2]\d)\b', text)]
    plausible = [y for y in years if 1980 <= y <= 2030]
    if plausible:
        return str(Counter(plausible).most_common(1)[0][0])
    return ""


# ---------------------------------------------------------------------------
# Categoria e nome sugeridos
# ---------------------------------------------------------------------------

def suggest_category(doc_type: str, text: str, filepath: str = "") -> str:
    dt = doc_type.lower()
    text_up = text.upper() if text else ""

    # Tipos de saúde → sempre SAUDE
    if doc_type in _SAUDE_TIPOS:
        return 'SAUDE'

    # Pasta /SAUDE/ + não é tipo explicitamente acadêmico → SAUDE
    if _in_saude_path(filepath) and doc_type not in (
        'Artigo Científico', 'Dissertação', 'Tese', 'TCC', 'Monografia', 'Trabalho Acadêmico'
    ):
        return 'SAUDE'

    # Tipos financeiros → sempre FINANCEIRO
    if doc_type in _FINANCEIRO_TIPOS:
        return 'FINANCEIRO'

    # Pasta /FINANCEIRO/ + não é tipo explicitamente acadêmico → FINANCEIRO
    _TIPOS_ACADEMICOS = {'Artigo Científico', 'Dissertação', 'Tese', 'TCC', 'Monografia', 'Trabalho Acadêmico'}
    if _in_financeiro_path(filepath) and doc_type not in _TIPOS_ACADEMICOS:
        return 'FINANCEIRO'

    FINANCEIRO = {'boleto', 'fatura', 'comprovante de matrícula', 'comprovante de matricula',
                  'comprovante', 'extrato', 'holerite'}
    JURIDICO = {'contrato', 'petição', 'peticao', 'decisão', 'decisao', 'jurídico'}
    PESSOAL = {'diploma', 'certificado', 'currículo', 'curriculo', 'histórico', 'historico'}
    ACADEMICO = {
        'artigo científico', 'dissertação', 'tese', 'monografia', 'tcc',
        'trabalho acadêmico', 'aula', 'avaliação', 'exercício', 'prova',
        'guia', 'apostila', 'apresentação', 'resumo', 'roteiro', 'edital',
        'projeto', 'requerimento', 'formulário', 'ficha',
    }
    if any(t in dt for t in FINANCEIRO):
        return 'FINANCEIRO'
    if any(kw in text_up for kw in ['BOLETO BANCÁRIO', 'BOLETO BANCARIO', 'NOSSO NÚMERO', 'HOLERITE']):
        return 'FINANCEIRO'
    if any(t in dt for t in JURIDICO):
        return 'JURIDICO'
    if any(t in dt for t in PESSOAL):
        return 'PESSOAL'
    if any(t in dt for t in ACADEMICO):
        return 'ACADEMICO'
    return 'GERAL'


def _safe(s: str, max_len: int = 120) -> str:
    return re.sub(r'[/\\:*?"<>|]', '-', s)[:max_len].strip()


def build_suggested_name(
    fn_info: Dict,
    doc_type: str,
    title: str,
    author: str,
    year: str,
    sobrenome: str = "",
) -> str:
    stem = fn_info['stem']
    ext = fn_info['ext']
    is_opaque = fn_info['is_opaque']

    # Higiene inicial dos metadados
    title = _clean_metadata_string(title)
    author = _clean_metadata_string(author)

    def extras() -> str:
        parts = [p for p in (author, year) if p]
        return (' - ' + ' - '.join(parts)) if parts else ''

    # 1. Padrões de slots fixos para Tipos Financeiros Específicos
    # Força a reconstrução para estes tipos se houver metadados mínimos
    if doc_type == 'Fatura':
        parts = ['Fatura']
        if title:  # Operadora
            parts.append(title)
        if author:  # Titular
            parts.append(author)
        if year:  # Data
            parts.append(year)
        if len(parts) >= 2: # No mínimo Fatura - Operadora
            return ' - '.join(parts) + ext

    if doc_type in ('Declaração de IRPF', 'DARF'):
        prefix = 'IRPF' if doc_type == 'Declaração de IRPF' else 'DARF'
        parts = [prefix]
        if author:  # Contribuinte
            parts.append(author)
        if year:  # Ano
            parts.append(year)
        if len(parts) >= 2: # No mínimo IRPF - Contribuinte ou IRPF - Ano
            return ' - '.join(parts) + ext

    # 2. Fallback para outros tipos financeiros
    if doc_type in _FINANCEIRO_TIPOS:
        if is_opaque:
            parts = [doc_type]
            if title:
                parts.append(title)
            if author:
                parts.append(author)
            if year:
                parts.append(year)
            return ' - '.join(parts) + ext
        # Se não for opaco, tentamos limpar o stem se ele for muito parecido com o tipo
        clean_stem = stem
        if doc_type.lower() in stem.lower() and len(stem) < len(doc_type) + 5:
            # ex: "Fatura.pdf" ou "DARF_2019.pdf" -> prefere reconstruir
            pass 
        else:
            return f"{stem}{extras()}{ext}"

    # 3. Tipos de saúde: opaque → usa tipo como prefixo; non-opaque → mantém stem + extras
    if doc_type in _SAUDE_TIPOS:
        if is_opaque:
            parts = [doc_type]
            if title:
                parts.append(title)
            if author:
                parts.append(author)
            if year:
                parts.append(year)
            return ' - '.join(parts) + ext
        return f"{stem}{extras()}{ext}"

    # 4. Trabalho Acadêmico com título: sempre usa template título-primeiro
    if doc_type == 'Trabalho Acadêmico' and title:
        parts = ['Trabalho', title]
        if author:
            parts.append(author)
        if year:
            parts.append(year)
        return ' - '.join(parts) + ext

    if not is_opaque:
        return f"{stem}{extras()}{ext}"

    # Nome opaco → reconstruir a partir de metadados
    tipo = doc_type or 'Documento'

    if 'boleto' in tipo.lower():
        return f"Boleto{extras()}{ext}" if (author or year) else f"Boleto - {_safe(stem)}{ext}"

    if doc_type == 'Artigo Científico':
        parts = ['Artigo']
        if title:
            parts.append(title)
        if sobrenome:
            parts.append(sobrenome)
        elif author:
            parts.append(author)
        if year:
            parts.append(year)
        return ' - '.join(parts) + ext

    titulo_part = title if title else None
    parts = [tipo] + ([titulo_part] if titulo_part else []) + ([author] if author else []) + ([year] if year else [])
    return ' - '.join(parts) + ext


# ---------------------------------------------------------------------------
# Debug: salva texto de PDFs com texto mas metadados incompletos
# ---------------------------------------------------------------------------

def _sanitize_for_filename(name: str) -> str:
    return re.sub(r'[^\w\-]', '_', name)[:60]


def save_debug_text(filepath: str, text: str) -> str:
    os.makedirs(DEBUG_DIR, exist_ok=True)
    base = _sanitize_for_filename(os.path.splitext(os.path.basename(filepath))[0])
    debug_path = os.path.join(DEBUG_DIR, f"{base}-page1.txt")
    with open(debug_path, 'w', encoding='utf-8') as f:
        f.write(text[:2000])
    return debug_path


# ---------------------------------------------------------------------------
# Análise de arquivo
# ---------------------------------------------------------------------------

def analyze_file(filepath: str) -> Dict:
    filename = os.path.basename(filepath)
    fn_info = analyze_filename(filename)
    ext = fn_info['ext'].lower()

    # Extração de texto: PDFs usam pipeline com OCR
    ocr_applied = False
    if ext == '.pdf':
        text, ocr_applied = _extract_pdf(filepath)
    else:
        text = extract_text(filepath)

    text_len = len(text.strip())

    extraction_reason: Optional[str] = None

    doc_type = fn_info['tipo_from_name']
    title = ""
    author = ""
    year = ""
    sobrenome = ""
    debug_path: Optional[str] = None
    nomes_relevantes: List[str] = []

    if text_len == 0:
        extraction_reason = "sem_texto"
        doc_type = doc_type or 'Desconhecido'
    else:
        nomes_relevantes = detecta_nomes_relevantes(text)
        in_saude = _in_saude_path(filepath)

        # 1. Detector de artigo de periódico (ISSN / citação) — roda mesmo em SAUDE
        artigo_meta: Optional[Dict] = None
        if ext == '.pdf':
            artigo_meta = detecta_artigo_academico(text)

        if artigo_meta:
            doc_type = 'Artigo Científico'
            title = artigo_meta['titulo']
            author = ', '.join(artigo_meta['autores']) if artigo_meta['autores'] else ""
            year = artigo_meta['ano']
            sobrenome = artigo_meta['autor_principal_sobrenome']

        elif in_saude or _is_saude_evidence(text):
            # 2. Contexto SAUDE: short-circuit — não chama detecta_capa_academica
            saude_meta = detecta_documento_saude(text, filepath)
            if saude_meta:
                # Preserva tipo do nome do arquivo se já é tipo de saúde específico
                if fn_info['tipo_from_name'] not in _SAUDE_TIPOS:
                    doc_type = saude_meta['tipo_documento']
                title = saude_meta['titulo']
                author = saude_meta['autor_principal']
                year = saude_meta['ano']
            else:
                type_from_content = detect_type_from_content(text)
                if fn_info['tipo_from_name']:
                    doc_type = fn_info['tipo_from_name']
                    if fn_info['tipo_from_name'] in ('Comprovante', 'Documento') and type_from_content:
                        doc_type = type_from_content
                else:
                    doc_type = type_from_content or 'Desconhecido'
                title = extract_title_from_content(text, doc_type)
                author = extract_author_from_content(text)
                year = extract_year_from_content(text)

        elif _in_financeiro_path(filepath) or _is_financeiro_evidence(text):
            # 3. Contexto FINANCEIRO: short-circuit — não chama detecta_capa_academica
            fin_meta = detecta_documento_financeiro(text, filepath)
            if fin_meta:
                # O tipo detectado no conteúdo tem prioridade se for específico, 
                # ou se o tipo do nome for muito genérico (Boleto, Comprovante)
                if fin_meta['tipo_documento'] not in ('Documento Financeiro', 'Desconhecido'):
                    doc_type = fin_meta['tipo_documento']
                elif not fn_info['tipo_from_name'] or fn_info['tipo_from_name'] in ('Boleto', 'Comprovante', 'Documento'):
                    doc_type = fin_meta['tipo_documento']
                
                title = fin_meta['titulo']
                author = fin_meta['autor_principal']
                year = fin_meta['ano']
            else:
                type_from_content = detect_type_from_content(text)
                if fn_info['tipo_from_name']:
                    doc_type = fn_info['tipo_from_name']
                    if fn_info['tipo_from_name'] in ('Comprovante', 'Documento') and type_from_content:
                        doc_type = type_from_content
                else:
                    doc_type = type_from_content or 'Desconhecido'
                title = extract_title_from_content(text, doc_type)
                author = extract_author_from_content(text)
                year = extract_year_from_content(text)

        else:
            # 4. Fluxo normal: tenta capa acadêmica → genérico
            capa_meta: Optional[Dict] = None
            if ext in ('.pdf', '.docx', '.doc', '.odt'):
                capa_meta = detecta_capa_academica(text)

            if capa_meta:
                doc_type = capa_meta['tipo_documento']
                title = capa_meta['titulo']
                author = capa_meta['autor_principal']
                year = capa_meta['ano']
            else:
                type_from_content = detect_type_from_content(text)
                if fn_info['tipo_from_name']:
                    doc_type = fn_info['tipo_from_name']
                    if fn_info['tipo_from_name'] in ('Comprovante', 'Documento') and type_from_content:
                        doc_type = type_from_content
                else:
                    doc_type = type_from_content or 'Desconhecido'
                title = extract_title_from_content(text, doc_type)
                author = extract_author_from_content(text)
                year = extract_year_from_content(text)

        has_metadata = bool(title or author or year)
        if not has_metadata:
            extraction_reason = "texto_ok_metadados_incompletos"
            if text_len >= 200 and ext in ('.pdf', '.docx', '.doc'):
                debug_path = save_debug_text(filepath, text)

    extraction_incomplete = extraction_reason is not None
    is_opaque = fn_info['is_opaque']

    if is_opaque and extraction_reason == "sem_texto":
        suggested_name = '[NÃO SUGERIDO — sem texto/OCR]'
    elif is_opaque and extraction_reason == "texto_ok_metadados_incompletos":
        suggested_name = '[NÃO SUGERIDO — texto lido, metadados não identificados]'
    else:
        suggested_name = build_suggested_name(fn_info, doc_type, title, author, year, sobrenome)

    return {
        'path': filepath,
        'filename': filename,
        'ext': ext,
        'is_opaque_name': is_opaque,
        'type': doc_type,
        'title': title,
        'author': author,
        'year': year,
        'text_len': text_len,
        'extraction_incomplete': extraction_incomplete,
        'extraction_reason': extraction_reason,
        'suggested_name': suggested_name,
        'suggested_category': suggest_category(doc_type, text, filepath),
        'nomes_relevantes': nomes_relevantes,
        'debug_path': debug_path,
        'ocr_aplicado': ocr_applied,
    }


# ---------------------------------------------------------------------------
# Relatório Markdown
# ---------------------------------------------------------------------------

def generate_report(results: list, report_path: str, base_dir: str = BASE_DIR) -> str:
    lines = [
        f'# V6 (piloto) — Propostas de renomeação/movimentação — {base_dir}\n',
        f'Total de arquivos analisados: {len(results)}\n',
        '---\n',
    ]

    for r in results:
        incomplete_str = 'false'
        if r['extraction_incomplete']:
            reason = r.get('extraction_reason', '')
            if reason == 'sem_texto':
                incomplete_str = 'true (sem texto/OCR)'
            elif reason == 'texto_ok_metadados_incompletos':
                incomplete_str = 'true (texto lido, metadados não identificados)'
            else:
                incomplete_str = 'true'

        lines.append(f"- [ ] {r['path']}  ")
        lines.append(f"      → SUGESTÃO DE NOME: {r['suggested_name']}  ")
        lines.append(f"      categoria_sugerida: {r['suggested_category']}  ")

        detail = (
            f"      (tipo: {r['type']}, "
            f"autor: {r['author'] or '—'}, "
            f"ano: {r['year'] or '—'}, "
            f"extração_incompleta: {incomplete_str}"
        )
        nomes = r.get('nomes_relevantes') or []
        if nomes:
            detail += f", nomes_relevantes: [{', '.join(nomes)}]"
        if r.get('ocr_aplicado'):
            detail += ", ocr_aplicado: true"
        if r.get('debug_path'):
            detail += f", debug: {r['debug_path']}"
        detail += ")  "
        lines.append(detail)
        lines.append('')

    sem_texto = sum(1 for r in results if r.get('extraction_reason') == 'sem_texto')
    texto_sem_meta = sum(1 for r in results if r.get('extraction_reason') == 'texto_ok_metadados_incompletos')
    complete = sum(1 for r in results if not r['extraction_incomplete'])

    lines += [
        '\n---\n',
        '## Estatísticas\n',
        f'- Total: {len(results)}',
        f'- Extração completa (metadados identificados): {complete}',
        f'- Sem texto / scan sem OCR: {sem_texto}',
        f'- Texto lido, metadados incompletos: {texto_sem_meta}',
        '',
        '### Por categoria sugerida:\n',
    ]
    for cat, count in sorted(Counter(r['suggested_category'] for r in results).items()):
        lines.append(f'- {cat}: {count}')

    lines.append('\n### Por tipo detectado:\n')
    for t, count in sorted(Counter(r['type'] for r in results).items()):
        lines.append(f'- {t}: {count}')

    os.makedirs(os.path.dirname(report_path), exist_ok=True)
    with open(report_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))
    return report_path


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

BUILD_TAG = "V6-PILOTO-BUILD-2026-04-28-19h30"


def main() -> None:
    import argparse
    print(f"DEBUG V6 PILOTO — {BUILD_TAG}")
    print(f"OCR module: {'disponível' if HAS_OCR_MODULE else 'não disponível (PyPDF2 apenas)'}")
    parser = argparse.ArgumentParser(description='V6 Vault Saneamento Piloto — Modo Leitura/Sugestão')
    parser.add_argument(
        '--base-dir',
        default=BASE_DIR,
        help=f'Diretório base para análise (padrão: {BASE_DIR})',
    )
    parser.add_argument(
        '--report',
        default=REPORT_PATH,
        help=f'Caminho do relatório Markdown a gerar (padrão: {REPORT_PATH})',
    )
    args = parser.parse_args()

    base_dir = os.path.expanduser(args.base_dir)
    report_path = os.path.expanduser(args.report)

    print('V6 Vault Saneamento Piloto — Modo Leitura/Sugestão')
    print(f'Diretório base: {base_dir}')
    print(f'Relatório:      {report_path}')
    print()

    if not os.path.isdir(base_dir):
        print(f'ERRO: diretório não encontrado: {base_dir}')
        sys.exit(1)

    files = sorted(
        os.path.join(root, f)
        for root, _, filenames in os.walk(base_dir)
        for f in filenames
        if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS
    )
    print(f'Arquivos encontrados: {len(files)}\n')

    results = []
    for i, fpath in enumerate(files, 1):
        name = os.path.basename(fpath)
        print(f'[{i:3}/{len(files)}] {name}')
        try:
            r = analyze_file(fpath)
            results.append(r)
            if i <= 10 or '5014-12955-1-PB' in name:
                ocr_flag = ' [OCR]' if r.get('ocr_aplicado') else ''
                print(f'  DEBUG EXTRACAO: ext={r["ext"]} texto_len={r["text_len"]} '
                      f'razao={r["extraction_reason"]}{ocr_flag}')
        except Exception as e:
            print(f'  ERRO: {e}')
            results.append({
                'path': fpath,
                'filename': name,
                'ext': os.path.splitext(fpath)[1].lower(),
                'is_opaque_name': True,
                'type': 'Erro',
                'title': '',
                'author': '',
                'year': '',
                'text_len': 0,
                'extraction_incomplete': True,
                'extraction_reason': 'erro',
                'suggested_name': '[ERRO NA ANÁLISE]',
                'suggested_category': 'GERAL',
                'nomes_relevantes': [],
                'debug_path': None,
                'ocr_aplicado': False,
            })

    print()
    report_out = generate_report(results, report_path, base_dir)

    sem_texto = sum(1 for r in results if r.get('extraction_reason') == 'sem_texto')
    texto_sem_meta = sum(1 for r in results if r.get('extraction_reason') == 'texto_ok_metadados_incompletos')
    complete = sum(1 for r in results if not r['extraction_incomplete'])
    n_ocr = sum(1 for r in results if r.get('ocr_aplicado'))

    print('=== RESUMO ===')
    print(f'Total de arquivos:              {len(results)}')
    print(f'Extração completa:              {complete}')
    print(f'Sem texto / scan sem OCR:       {sem_texto}')
    print(f'Texto lido, metadados incompletos: {texto_sem_meta}')
    print(f'OCR aplicado (tesseract):       {n_ocr}')
    if HAS_OCR_MODULE:
        stats = _get_ocr_stats()
        print(f'  PyPDF2:      {stats.get("pypdf2", 0)}  |  '
              f'pdftotext: {stats.get("pdftotext", 0)}  |  '
              f'tesseract: {stats.get("tesseract", 0)}  |  '
              f'skipped: {stats.get("tesseract_skipped", 0)}')
    print(f'Relatório gerado em:            {report_out}')
    print()

    # Resumo específico SAUDE
    saude_results = [r for r in results if '/SAUDE/' in r['path'].replace('\\', '/').upper()]
    if saude_results:
        saude_completo = sum(1 for r in saude_results if not r['extraction_incomplete'])
        saude_sem_texto = sum(1 for r in saude_results if r.get('extraction_reason') == 'sem_texto')
        saude_sem_meta = sum(1 for r in saude_results if r.get('extraction_reason') == 'texto_ok_metadados_incompletos')
        saude_ocr = sum(1 for r in saude_results if r.get('ocr_aplicado'))
        saude_tipos = Counter(r['type'] for r in saude_results)

        print('=== RESUMO SAUDE ===')
        print(f'Total SAUDE:                    {len(saude_results)}')
        print(f'SAUDE - extração completa:      {saude_completo}')
        print(f'SAUDE - sem texto/OCR:          {saude_sem_texto}')
        print(f'SAUDE - texto_ok_metadados_incompletos: {saude_sem_meta}')
        print(f'SAUDE - OCR aplicado:           {saude_ocr}')
        print('SAUDE - tipos:')
        for tipo in ['Resultado de Exame', 'Laudo', 'Guia de Exame', 'Questionário Clínico',
                     'Receita', 'Documento de Saúde', 'Trabalho Acadêmico', 'Desconhecido']:
            n = saude_tipos.get(tipo, 0)
            if n:
                print(f'  {tipo}: {n}')
        outros = {k: v for k, v in saude_tipos.items()
                  if k not in ('Resultado de Exame', 'Laudo', 'Guia de Exame', 'Questionário Clínico',
                                'Receita', 'Documento de Saúde', 'Trabalho Acadêmico', 'Desconhecido')}
        for tipo, n in sorted(outros.items()):
            print(f'  {tipo}: {n}')
        print()

    # Resumo específico FINANCEIRO
    fin_results = [r for r in results if '/FINANCEIRO/' in r['path'].replace('\\', '/').upper()]
    if fin_results:
        fin_completo = sum(1 for r in fin_results if not r['extraction_incomplete'])
        fin_sem_texto = sum(1 for r in fin_results if r.get('extraction_reason') == 'sem_texto')
        fin_sem_meta = sum(1 for r in fin_results if r.get('extraction_reason') == 'texto_ok_metadados_incompletos')
        fin_ocr = sum(1 for r in fin_results if r.get('ocr_aplicado'))
        fin_trab_acad = sum(1 for r in fin_results if r.get('type') == 'Trabalho Acadêmico')
        fin_tipos = Counter(r['type'] for r in fin_results)

        print('=== RESUMO FINANCEIRO ===')
        print(f'Total FINANCEIRO:               {len(fin_results)}')
        print(f'FINANCEIRO - extração completa: {fin_completo}')
        print(f'FINANCEIRO - sem texto/OCR:     {fin_sem_texto}')
        print(f'FINANCEIRO - texto_ok_metadados_incompletos: {fin_sem_meta}')
        print(f'FINANCEIRO - OCR aplicado:      {fin_ocr}')
        print(f'FINANCEIRO - Trabalho Acadêmico (falso positivo): {fin_trab_acad}')
        print('FINANCEIRO - tipos:')
        for tipo in ['Nota Fiscal', 'Comprovante de Transferência', 'Recibo de Salário/Holerite',
                     'Extrato Bancário', 'Declaração de IRPF', 'Fatura', 'Boleto',
                     'Documento Financeiro', 'Trabalho Acadêmico', 'Desconhecido']:
            n = fin_tipos.get(tipo, 0)
            if n:
                print(f'  {tipo}: {n}')
        outros_fin = {k: v for k, v in fin_tipos.items()
                      if k not in ('Nota Fiscal', 'Comprovante de Transferência',
                                   'Recibo de Salário/Holerite', 'Extrato Bancário',
                                   'Declaração de IRPF', 'Fatura', 'Boleto',
                                   'Documento Financeiro', 'Trabalho Acadêmico', 'Desconhecido')}
        for tipo, n in sorted(outros_fin.items()):
            print(f'  {tipo}: {n}')
        print()

    print('=== GOLDEN TEST ===')
    golden = next((r for r in results if '5014-12955-1-PB' in r['path']), None)
    if golden:
        print(f'  {golden["filename"]}')
        print(f'  → {golden["suggested_name"]}')
        print(f'  [tipo: {golden["type"]}, autor: {golden["author"]}, ano: {golden["year"]}]')
        print(f'  incompleto: {golden["extraction_incomplete"]} / razão: {golden["extraction_reason"]}')
    else:
        print('  (arquivo não encontrado no diretório de teste)')
    print()
    print('=== GOLDEN CAPA ACADÊMICA ===')
    golden_capa = next((r for r in results if r['filename'] == 'ArtigoAV2.docx'), None)
    if golden_capa:
        print(f'  {golden_capa["filename"]}')
        print(f'  → {golden_capa["suggested_name"]}')
        print(f'  [tipo: {golden_capa["type"]}, categoria: {golden_capa["suggested_category"]}, '
              f'autor: {golden_capa["author"]}, ano: {golden_capa["year"]}]')
        print(f'  nomes_relevantes: {golden_capa.get("nomes_relevantes", [])}')
        ok_tipo = golden_capa['type'] != 'Desconhecido'
        ok_cat  = golden_capa['suggested_category'] == 'ACADEMICO'
        ok_nome = 'Trabalho' in golden_capa['suggested_name'] and golden_capa['title'] in golden_capa['suggested_name']
        print(f'  PASS tipo≠Desconhecido: {ok_tipo} | PASS cat=ACADEMICO: {ok_cat} | PASS título-no-nome: {ok_nome}')
    else:
        print('  (ArtigoAV2.docx não encontrado no diretório)')
    print()
    print('=== EXEMPLOS (5 primeiros) ===')
    for r in results[:5]:
        print(f'  {r["filename"]}')
        print(f'    → {r["suggested_name"]}')
        print(f'    [{r["type"]} | {r["suggested_category"]} | razão={r["extraction_reason"]}]')


if __name__ == '__main__':
    main()
